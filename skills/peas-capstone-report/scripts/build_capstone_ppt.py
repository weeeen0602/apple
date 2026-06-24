#!/usr/bin/env python3
"""Build report/專題報告.pptx from report/專題報告.md."""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx required. Run: uv add python-pptx", file=sys.stderr)
    raise SystemExit(1)


IMAGE_RE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
HEADING_RE = re.compile(r"^(#{1,3})\s+(.+)$")
BULLET_RE = re.compile(r"^[-*]\s+(.+)$")


def parse_md(md_text: str) -> dict:
    lines = md_text.splitlines()
    title = "專題報告"
    meta: list[str] = []
    sections: list[tuple[str, list[str], list[str]]] = []
    current_title = ""
    current_bullets: list[str] = []
    current_images: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_bullets, current_images
        if current_title:
            sections.append((current_title, current_bullets.copy(), current_images.copy()))
        current_title = ""
        current_bullets = []
        current_images = []

    for line in lines:
        if line.startswith("# ") and not sections and not current_title:
            title = line[2:].strip()
            continue
        if line.startswith("**組別") or line.startswith("**日期"):
            meta.append(line.strip("* "))
            continue
        m = HEADING_RE.match(line)
        if m:
            flush()
            current_title = m.group(2).strip()
            continue
        img = IMAGE_RE.search(line)
        if img:
            current_images.append(img.group(1).strip())
            continue
        b = BULLET_RE.match(line)
        if b and current_title:
            current_bullets.append(b.group(1).strip())
            continue
        if line.strip() and current_title and not line.startswith("---"):
            if not line.startswith("!["):
                current_bullets.append(line.strip())

    flush()
    return {"title": title, "meta": meta, "sections": sections}


def add_title_slide(prs: Presentation, title: str, meta: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "\n".join(meta) if meta else "Agent Studio 專題報告"


def _content_layout(prs: Presentation):
    for idx in (1, 5, 6):
        if idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[idx]
            if len(layout.placeholders) > 1:
                return layout
    return prs.slide_layouts[1]


def add_content_slide(
    prs: Presentation,
    heading: str,
    bullets: list[str],
    report_dir: Path,
    images: list[str],
) -> None:
    slide = prs.slides.add_slide(_content_layout(prs))
    slide.shapes.title.text = heading
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets[:8]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

    for img_rel in images:
        img_path = (report_dir / img_rel).resolve()
        if img_path.is_file():
            slide.shapes.add_picture(
                str(img_path),
                Inches(0.8),
                Inches(2.2),
                width=Inches(8.5),
            )


def add_image_slide(prs: Presentation, heading: str, img_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0])
    slide.shapes.title.text = heading
    if img_path.is_file():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0.8),
            Inches(1.6),
            height=Inches(5.0),
        )


def build(report_dir: Path) -> Path:
    md_path = report_dir / "專題報告.md"
    if not md_path.is_file():
        raise FileNotFoundError(f"Missing {md_path}")

    data = parse_md(md_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, data["title"], data["meta"])

    for heading, bullets, images in data["sections"]:
        if "附錄" in heading:
            continue
        if images and len(bullets) <= 2:
            add_image_slide(prs, heading, (report_dir / images[0]).resolve())
            for b in bullets:
                slide = prs.slides.add_slide(_content_layout(prs))
                slide.shapes.title.text = heading
                slide.placeholders[1].text = b
        else:
            add_content_slide(prs, heading, bullets, report_dir, images)

    assets = report_dir / "assets"
    demo_images = sorted(assets.glob("demo-*.png"))
    for i, demo in enumerate(demo_images, start=1):
        add_image_slide(prs, f"Demo {i}", demo)

    out = report_dir / "專題報告.pptx"
    prs.save(str(out))
    return out


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--report-dir", type=Path, default=Path("report"))
    args = parser.parse_args()
    try:
        out = build(args.report_dir)
    except FileNotFoundError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    print(f"OK: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
